import streamlit as st
import os
import base64
import json
import re
from pptx import Presentation
from docx import Document as DocxDocument
from mistralai import Mistral
class RecursiveCharacterTextSplitter:
    def __init__(self, chunk_size=500, chunk_overlap=100, separators=None, length_function=len):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    def split_text(self, text):
        chunks = []
        start = 0
        while start < len(text):
            end = start + self.chunk_size
            chunks.append(text[start:end])
            start += self.chunk_size - self.chunk_overlap
        return [c for c in chunks if c.strip()]
from typing import Optional, List, Dict, Tuple
import time
from io import BytesIO
from datetime import datetime
import hashlib


# PDF Libraries
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("⚠️ ReportLab not installed. PDF export will be disabled.")
    print("Install with: pip install reportlab")

# ======================================
# CONFIGURATION
# ======================================
st.set_page_config(
    layout="wide",
    page_title="AI Study Assistant",
    page_icon="🎓",
    initial_sidebar_state="expanded"
)

# ======================================
# SIMPLE TRANSLATION HELPER
# ======================================
def T(ar: str, en: str) -> str:
    lang = st.session_state.get('ui_lang', 'ar')
    return ar if lang == 'ar' else en

# ======================================
# ADVANCED CACHING SYSTEM - VERY FAST
# ======================================
class FastCache:
    """نظام كاش متطور للغاية"""
    
    @staticmethod
    @st.cache_resource
    def get_mistral_client(api_key: str):
        """Cache Mistral client forever"""
        try:
            return Mistral(api_key=api_key)
        except Exception as e:
            st.error(f"❌ Failed to initialize Mistral client: {e}")
            return None
    
    @staticmethod
    @st.cache_data(ttl=3600, show_spinner=False)
    def get_file_hash(file_content: bytes) -> str:
        """Get file hash for caching"""
        return hashlib.md5(file_content).hexdigest()
    
    @staticmethod
    @st.cache_data(ttl=3600, show_spinner=False)
    def cache_extracted_text(file_hash: str, file_name: str, file_type: str, file_content: bytes) -> Tuple[str, dict]:
        """Cache extracted text based on file hash"""
        # This will be populated by the actual extraction function
        # We're just creating a placeholder for the cache key
        return "", {}
    
    @staticmethod
    @st.cache_data(ttl=7200, show_spinner=False)
    def cache_generated_features(text_hash: str, text_length: int) -> Dict:
        """Cache generated features - VERY IMPORTANT for speed"""
        # This is just a cache key placeholder
        return {}

# Initialize cache
cache = FastCache()

def get_text_hash(text: str) -> str:
    """Generate a hash for the text to use as cache key."""
    return hashlib.md5(text.encode('utf-8')).hexdigest()

# ======================================
# CHATBOT FUNCTIONS (أسماء)
# ======================================
def initialize_chat_session(filename: str, extracted_text: str):
    """Initialize a new chat session for a specific file."""
    if 'chat_sessions' not in st.session_state:
        st.session_state['chat_sessions'] = {}

    if filename not in st.session_state['chat_sessions']:
        st.session_state['chat_sessions'][filename] = {
            'messages': [],
            'context': extracted_text[:10000],
            'created_at': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            'total_tokens_used': 0
        }

    return st.session_state['chat_sessions'][filename]


def get_chat_context(filename: str) -> str:
    """Get the context text for the chat session."""
    if 'chat_sessions' in st.session_state and filename in st.session_state['chat_sessions']:
        return st.session_state['chat_sessions'][filename]['context']
    return ""


def add_chat_message(filename: str, role: str, content: str, tokens: int = 0):
    """Add a message to the chat history."""
    if 'chat_sessions' in st.session_state and filename in st.session_state['chat_sessions']:
        st.session_state['chat_sessions'][filename]['messages'].append({
            'role': role,
            'content': content,
            'timestamp': datetime.now().strftime("%H:%M:%S")
        })
        st.session_state['chat_sessions'][filename]['total_tokens_used'] += tokens

        if len(st.session_state['chat_sessions'][filename]['messages']) > 20:
            st.session_state['chat_sessions'][filename]['messages'] = \
                st.session_state['chat_sessions'][filename]['messages'][-20:]


def get_chat_messages(filename: str) -> List[Dict]:
    """Get chat messages for a specific file."""
    if 'chat_sessions' in st.session_state and filename in st.session_state['chat_sessions']:
        return st.session_state['chat_sessions'][filename]['messages']
    return []


def clear_chat_session(filename: str):
    """Clear chat history for a specific file."""
    if 'chat_sessions' in st.session_state and filename in st.session_state['chat_sessions']:
        st.session_state['chat_sessions'][filename]['messages'] = []
        st.session_state['chat_sessions'][filename]['total_tokens_used'] = 0
        st.session_state[f'chat_reset_{filename}'] = True


def format_error_message(error: Exception, context: str = "") -> str:
    """Format error messages based on error type."""
    error_str = str(error).lower()
    if "rate limit" in error_str or "429" in error_str:
        return "⏳ API rate limit reached. Please wait a moment and try again."
    elif "authentication" in error_str or "401" in error_str or "api key" in error_str:
        return "🔑 Invalid or expired API key. Please check your Mistral API key."
    elif "timeout" in error_str:
        return "⏱️ Request timeout. The file might be too complex or the server is busy."
    elif "connection" in error_str:
        return "🌐 Connection error. Please check your internet connection."
    else:
        return f"❌ {context}: {str(error)}"


def generate_chat_response(client, filename: str, user_message: str) -> str:
    """Generate AI response for chat based on document context."""
    try:
        chat_session = st.session_state['chat_sessions'][filename]
        context_text = chat_session['context']

        conversation_history = []
        for msg in chat_session['messages'][-10:]:
            conversation_history.append(f"{msg['role'].upper()}: {msg['content']}")

        history_text = "\n".join(conversation_history[-5:]) if conversation_history else "No previous conversation."

        system_prompt = f"""You are an AI Study Assistant specialized in explaining and discussing the uploaded document.

DOCUMENT CONTEXT:
{context_text[:8000]}

IMPORTANT INSTRUCTIONS:
1. You MUST answer based ONLY on the document content provided above
2. If the question is outside the document scope, politely say you can only answer based on the uploaded document
3. Support both English and Arabic languages seamlessly
4. Be helpful, educational, and precise
5. For Arabic questions, answer in Arabic. For English questions, answer in English
6. You can:
   - Explain concepts from the document
   - Answer specific questions about the content
   - Provide examples related to the material
   - Help with study strategies for this document
   - Clarify confusing parts
7. Keep answers concise but informative

PREVIOUS CONVERSATION:
{history_text}

USER'S QUESTION:
{user_message}

ANSWER:"""

        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message}
            ],
            temperature=0.7,
            max_tokens=1500
        )

        response_text = response.choices[0].message.content
        tokens_used = response.usage.total_tokens if hasattr(response, 'usage') else 0

        add_chat_message(filename, "user", user_message, tokens=0)
        add_chat_message(filename, "assistant", response_text, tokens=tokens_used)

        return response_text

    except Exception as e:
        error_msg = format_error_message(e, "Chat generation")
        return f"❌ Error: {error_msg}"


def export_chat_history(filename: str):
    """Export chat history to a PDF file."""
    if 'chat_sessions' in st.session_state and filename in st.session_state['chat_sessions']:
        chat_session = st.session_state['chat_sessions'][filename]
        messages = chat_session['messages']

        pdf_content = f"Chat History for: {filename}\n"
        pdf_content += f"Created: {chat_session['created_at']}\n"
        pdf_content += f"Total Tokens Used: {chat_session['total_tokens_used']}\n"
        pdf_content += "="*50 + "\n\n"

        for msg in messages:
            pdf_content += f"{msg['role'].upper()} ({msg.get('timestamp', '')}):\n"
            pdf_content += f"{msg['content']}\n"
            pdf_content += "-"*40 + "\n"

        if PDF_AVAILABLE:
            try:
                pdf_buffer = create_pdf_from_text(pdf_content, f"Chat History - {filename}")
                st.download_button(
                    label=T("📥 تحميل المحادثة (PDF)", "📥 Download Chat History (PDF)"),
                    data=pdf_buffer,
                    file_name=f"{filename}_chat_history.pdf",
                    mime="application/pdf",
                    key=f"export_chat_{filename}"
                )
            except Exception as e:
                st.error(f"❌ Error creating PDF: {str(e)[:100]}")
                st.download_button(
                    label=T("📥 تحميل المحادثة (TXT)", "📥 Download Chat History (TXT)"),
                    data=pdf_content,
                    file_name=f"{filename}_chat_history.txt",
                    mime="text/plain",
                    key=f"export_chat_txt_{filename}"
                )
        else:
            st.download_button(
                label=T("📥 تحميل المحادثة (TXT)", "📥 Download Chat History (TXT)"),
                data=pdf_content,
                file_name=f"{filename}_chat_history.txt",
                mime="text/plain",
                key=f"export_chat_txt_{filename}"
            )


def display_chat_interface(client, filename: str):
    """Display the chat interface for a specific document."""
    st.subheader(T("💬 مساعد المحادثة مع المستند", "💬 Document Chat Assistant"))
    st.markdown(T("**اتكلم مع المستند بتاعك:** ", "**Chat with your document:** ") + filename)
    st.markdown(T(
        "*تقدر تسأل بالعربي أو الإنجليزي*",
        "*You can ask questions in English or العربية*"
    ))
    st.markdown("---")

    if 'extracted_text' in st.session_state:
        initialize_chat_session(filename, st.session_state['extracted_text'])

    chat_container = st.container()
    with chat_container:
        messages = get_chat_messages(filename)
        if not messages:
            st.info(T(
                "💡 ابدأ واسأل أي سؤال عن المستند بالعربي أو الإنجليزي.",
                "💡 Start chatting! Ask questions about your document in English or Arabic."
            ))
            st.info(T(
                "📚 أمثلة: 'اشرح المفاهيم الرئيسية', 'ما هي النقاط المهمة؟', 'اديني أمثلة', 'Explain the main concepts'",
                "📚 Examples: 'Explain the main concepts', 'ما هي النقاط الرئيسية؟', 'Give me examples', 'اشرح لي هذا الجزء'"
            ))
        else:
            for msg in messages:
                if msg['role'] == 'user':
                    with st.chat_message("user", avatar="👤"):
                        st.markdown(msg['content'])
                        st.caption(f"⏰ {msg.get('timestamp', '')}")
                else:
                    with st.chat_message("assistant", avatar="🤖"):
                        st.markdown(msg['content'])
                        st.caption(f"⏰ {msg.get('timestamp', '')}")

        st.markdown("---")

        chat_input_key = f"chat_input_{hash(filename) % 10000}"
        
        # Use chat_input instead of form for better performance
        user_input = st.chat_input(
            T("اكتب سؤالك هنا...", "Type your message..."),
            key=chat_input_key
        )

        if user_input:
            if 'extracted_text' not in st.session_state:
                st.error(T(
                    "❌ لا يوجد مستند محمل. من فضلك ارفع مستند أولاً.",
                    "❌ No document loaded. Please upload a document first."
                ))
                st.stop()

            with st.spinner(T("🤖 جاري التفكير...", "🤖 Thinking...")):
                _ = generate_chat_response(client, filename, user_input)
            st.rerun()

    st.markdown("---")
    col1, col2, col3 = st.columns([1, 1, 2])

    with col1:
        if st.button(T("🗑️ مسح المحادثة", "🗑️ Clear Chat"), key=f"clear_{filename}", use_container_width=True):
            clear_chat_session(filename)
            st.rerun()

    with col2:
        if st.button(T("💾 حفظ المحادثة", "💾 Export Chat"), key=f"export_{filename}", use_container_width=True):
            export_chat_history(filename)

    with col3:
        if 'chat_sessions' in st.session_state and filename in st.session_state['chat_sessions']:
            token_count = st.session_state['chat_sessions'][filename]['total_tokens_used']
            st.caption(T(
                f"عدد التوكنز المستخدمة: {token_count:,}",
                f"Tokens used: {token_count:,}"
            ))

    st.markdown(T("### ⚡ إجراءات سريعة", "### ⚡ Quick Actions"))
    quick_cols = st.columns(4)
    quick_actions = {
        T("📋 تلخيص", "📋 Summarize"): "Please summarize the main points of this document in a clear and concise way.",
        T("❓ أسئلة وأجوبة", "❓ Ask Questions"): "Generate 5 important questions based on this document with their answers.",
        T("🔍 شرح", "🔍 Explain"): "Explain the most complex concept in this document in simple terms.",
        T("🌍 ترجمة النقاط الرئيسية", "🌍 Translate"): "Translate the key points of this document to Arabic."
    }

    for idx, (btn_text, action_text) in enumerate(quick_actions.items()):
        with quick_cols[idx]:
            if st.button(btn_text, key=f"quick_{idx}_{filename}", use_container_width=True):
                if 'extracted_text' in st.session_state:
                    with st.spinner(T("🤖 جاري التجهيز...", "🤖 Thinking...")):
                        _ = generate_chat_response(client, filename, action_text)
                    st.rerun()

# ======================================
# FILE VALIDATION
# ======================================
def validate_file(file, max_size_mb: float = 50) -> Tuple[bool, str]:
    """Validate file before processing."""
    try:
        if not file or file.size == 0:
            return False, "❌ Empty file detected"
        size_mb = file.size / (1024 * 1024)
        if size_mb > max_size_mb:
            return False, f"❌ File too large: {size_mb:.1f}MB (Max: {max_size_mb}MB)"
        return True, f"✅ Valid file ({size_mb:.2f}MB)"
    except Exception as e:
        return False, f"❌ Validation error: {str(e)}"

# ======================================
# TEXT CLEANING & PREPROCESSING
# ======================================
def clean_text(text: str) -> str:
    """Clean the extracted text - Basic cleaning."""
    if not text:
        return ""
    text = text.replace("\t", " ")
    text = text.replace("\xa0", " ")
    text = " ".join(text.split())
    return text.strip()


def chunk_text(text: str, chunk_size: int = 500, chunk_overlap: int = 100) -> List[str]:
    """Split text into overlapping chunks."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )
    chunks = splitter.split_text(text)
    return chunks

# ======================================
# HELPERS
# ======================================
def get_file_extension(filename: str) -> str:
    """Get file extension safely."""
    return os.path.splitext(filename)[1].lower()

# ======================================
# EXTRACTION FUNCTIONS (يوسف)
# ======================================
def extract_from_word(docx_file) -> Tuple[str, dict]:
    """Extract text from Word document."""
    try:
        doc = DocxDocument(docx_file)
        all_text = []
        metadata = {
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
            "sections": len(doc.sections),
            "errors": []
        }

        for para in doc.paragraphs:
            if para.text.strip():
                all_text.append(para.text)

        for table in doc.tables:
            all_text.append("\n[Table Start]")
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    all_text.append(" | ".join(row_text))
            all_text.append("[Table End]\n")

        final_text = "\n\n".join(all_text)
        return clean_text(final_text), metadata

    except Exception as e:
        error_msg = format_error_message(e, "Word extraction")
        return error_msg, {"error": str(e)}

def extract_from_txt(txt_file) -> Tuple[str, dict]:
    """Extract text from TXT file with multiple encoding support."""
    encodings = ['utf-8', 'utf-8-sig', 'cp1256', 'latin-1', 'iso-8859-1']
    for encoding in encodings:
        try:
            txt_file.seek(0)
            content = txt_file.read().decode(encoding)
            lines = content.split('\n')
            metadata = {
                "lines": len(lines),
                "characters": len(content),
                "words": len(content.split()),
                "encoding": encoding
            }
            return clean_text(content), metadata
        except (UnicodeDecodeError, AttributeError):
            continue

    return "Error: Could not decode text file with supported encodings", {
        "error": "Encoding not supported",
        "tried_encodings": encodings
    }

def pptx_extract_text(pptx_file, client) -> Tuple[str, dict]:
    """Extract text + tables + OCR from images inside PPTX."""
    try:
        prs = Presentation(pptx_file)
        all_slides_output = []
        metadata = {
            "total_slides": len(prs.slides),
            "slides_with_images": 0,
            "slides_with_tables": 0,
            "total_images_ocr": 0,
            "successful_ocr": 0,
            "failed_ocr": 0,
            "total_tables": 0,
            "errors": []
        }

        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_text_parts = [f"=== Slide {slide_number} ==="]
            has_images = False
            has_tables = False

            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            slide_text_parts.append(text)
                    elif hasattr(shape, "text") and isinstance(shape.text, str):
                        text = shape.text.strip()
                        if text:
                            slide_text_parts.append(text)
                except Exception as e:
                    metadata["errors"].append(f"Slide {slide_number}: Text extraction error - {str(e)}")

            for shape in slide.shapes:
                try:
                    if shape.has_table:
                        has_tables = True
                        metadata["total_tables"] += 1
                        table = shape.table
                        slide_text_parts.append("\n[Table Start]")
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_data:
                                slide_text_parts.append(" | ".join(row_data))
                        slide_text_parts.append("[Table End]\n")
                except Exception as e:
                    metadata["errors"].append(f"Slide {slide_number}: Table extraction error - {str(e)}")

            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "image"):
                    has_images = True
                    metadata["total_images_ocr"] += 1
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        mime = image.ext if image.ext in ["png", "jpg", "jpeg"] else "png"

                        image_size_mb = len(image_bytes) / (1024 * 1024)
                        if image_size_mb > 10:
                            metadata["errors"].append(
                                f"Slide {slide_number}, Image {shape_idx}: Image too large ({image_size_mb:.1f}MB)"
                            )
                            metadata["failed_ocr"] += 1
                            continue

                        encoded = base64.b64encode(image_bytes).decode("utf-8")
                        document = {
                            "type": "image_url",
                            "image_url": f"data:image/{mime};base64,{encoded}"
                        }

                        response = client.ocr.process(
                            model="mistral-ocr-latest",
                            document=document,
                            include_image_base64=False
                        )

                        if not hasattr(response, 'pages') or not response.pages:
                            metadata["failed_ocr"] += 1
                            continue

                        ocr_text_found = False
                        for page in response.pages:
                            page_text = None
                            if hasattr(page, 'markdown') and page.markdown:
                                page_text = page.markdown.strip()
                            elif hasattr(page, 'text') and page.text:
                                page_text = page.text.strip()

                            if page_text:
                                slide_text_parts.append(
                                    f"\n[Image OCR - Shape {shape_idx}]\n{page_text}"
                                )
                                ocr_text_found = True
                                metadata["successful_ocr"] += 1

                        if not ocr_text_found:
                            metadata["failed_ocr"] += 1
                    except Exception as e:
                        metadata["errors"].append(
                            f"Slide {slide_number}, Image {shape_idx}: {str(e)}"
                        )
                        metadata["failed_ocr"] += 1

            if has_images:
                metadata["slides_with_images"] += 1
            if has_tables:
                metadata["slides_with_tables"] += 1

            slide_output = "\n".join(slide_text_parts)
            all_slides_output.append(slide_output)

        final_text = "\n\n--- Slide Break ---\n\n".join(all_slides_output)
        return final_text.strip(), metadata

    except Exception as e:
        error_msg = format_error_message(e, "PPTX processing")
        return error_msg, {"error": str(e)}

def ocr_mistral(file, filename: str, client) -> Tuple[str, dict]:
    """OCR for images/PDF using Mistral."""
    metadata = {
        "filename": filename,
        "file_type": file.type,
        "pages_processed": 0,
        "errors": []
    }
    try:
        file_bytes = file.read()
        file.seek(0)
        file_size_mb = len(file_bytes) / (1024 * 1024)
        metadata["file_size_mb"] = round(file_size_mb, 2)

        if file_size_mb > 50:
            error_msg = f"❌ File exceeds Mistral OCR limit of 50MB (Current size: {file_size_mb:.1f}MB)"
            metadata["errors"].append(error_msg)
            return error_msg, metadata

        if len(file_bytes) == 0:
            return "Error: Empty file", metadata

        encoded = base64.b64encode(file_bytes).decode("utf-8")
        mime = file.type
        ext = get_file_extension(filename)

        if ext == ".pdf":
            document = {
                "type": "document_url",
                "document_url": f"data:application/pdf;base64,{encoded}"
            }
        else:
            document = {
                "type": "image_url",
                "image_url": f"data:{mime};base64:{encoded}"
            }

        response = client.ocr.process(
            model="mistral-ocr-latest",
            document=document,
            include_image_base64=False
        )

        final_text = ""
        pages = getattr(response, "pages", [])
        metadata["pages_processed"] = len(pages)

        for page in pages:
            page_text = getattr(page, "markdown", None) or getattr(page, "text", "")
            final_text += page_text + "\n\n"

        return final_text.strip(), metadata

    except Exception as e:
        error_msg = format_error_message(e, "OCR processing")
        metadata["errors"].append(error_msg)
        return error_msg, metadata

# ======================================
# PDF EXPORT FUNCTION - WITH ARABIC SUPPORT
# ======================================
def create_pdf_from_text(text: str, title: str) -> BytesIO:
    """Convert text to PDF with Arabic and English support"""
    if not PDF_AVAILABLE:
        raise Exception("ReportLab not installed. Install with: pip install reportlab")

    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from bidi.algorithm import get_display
        import arabic_reshaper
        ARABIC_SUPPORT = True
    except ImportError:
        ARABIC_SUPPORT = False
        print("⚠️ Arabic support disabled. Install: pip install python-bidi arabic-reshaper")

    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        topMargin=0.5*inch,
        bottomMargin=0.5*inch
    )

    font_name = 'Helvetica'
    if ARABIC_SUPPORT:
        try:
            arabic_fonts = [
                'C:/Windows/Fonts/arial.ttf',
                'C:/Windows/Fonts/tahoma.ttf',
                'C:/Windows/Fonts/calibri.ttf',
                '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
                '/System/Library/Fonts/Supplemental/Arial.ttf'
            ]
            for font_path in arabic_fonts:
                if os.path.exists(font_path):
                    pdfmetrics.registerFont(TTFont('Arabic', font_path))
                    font_name = 'Arabic'
                    break
        except Exception as e:
            print(f"⚠️ Could not load Arabic font: {e}")
            font_name = 'Helvetica'

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=20,
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName=font_name
    )
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=14,
        spaceAfter=12,
        spaceBefore=12,
        fontName=font_name
    )
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=11,
        leading=16,
        spaceAfter=8,
        fontName=font_name,
        alignment=TA_JUSTIFY
    )

    def contains_arabic(t):
        return any('\u0600' <= c <= '\u06FF' for c in t)

    def process_text(t):
        if ARABIC_SUPPORT and contains_arabic(t):
            try:
                reshaped = arabic_reshaper.reshape(t)
                return get_display(reshaped)
            except Exception:
                return t
        return t

    story = []

    try:
        processed_title = process_text(title)
        story.append(Paragraph(processed_title, title_style))
    except Exception:
        story.append(Paragraph(title, title_style))

    story.append(Spacer(1, 0.3*inch))

    lines = text.split('\n')
    for line in lines:
        line = line.strip()
        if not line:
            story.append(Spacer(1, 0.05*inch))
            continue

        line = line.replace('═', '-').replace('│', '|').replace('└', '+')
        line = line.replace('├', '+').replace('┌', '+').replace('─', '-')

        is_heading = (
            line.startswith('##') or
            line.startswith('###') or
            line.startswith('🎯') or
            line.startswith('📌') or
            (line.startswith('**') and line.endswith('**'))
        )

        line = line.replace('#', '').replace('**', '').strip()
        if not line:
            continue

        processed_line = process_text(line)

        try:
            if is_heading:
                story.append(Paragraph(processed_line, heading_style))
            elif line.startswith('-') or line.startswith('•') or line.startswith('+'):
                clean_line = line.lstrip('-•+').strip()
                processed_clean = process_text(clean_line)
                story.append(Paragraph(f"• {processed_clean}", normal_style))
            else:
                story.append(Paragraph(processed_line, normal_style))
        except Exception:
            try:
                story.append(Paragraph(line, normal_style))
            except Exception:
                continue

    try:
        doc.build(story)
    except Exception as e:
        print(f"⚠️ PDF build error: {e}")

    buffer.seek(0)
    return buffer

def create_download_button_pdf(content: str, base_filename: str, title: str, key_prefix: str):
    """Create download button for PDF format only"""
    if PDF_AVAILABLE:
        try:
            pdf_buffer = create_pdf_from_text(content, title)
            st.download_button(
                label=T("📕 تحميل PDF", "📕 Download PDF"),
                data=pdf_buffer,
                file_name=f"{base_filename}.pdf",
                mime="application/pdf",
                key=f"{key_prefix}_pdf",
                use_container_width=True
            )
        except Exception as e:
            st.warning(f"⚠️ PDF Error: {str(e)[:50]}")
            st.button("📕 PDF (Error)", disabled=True, use_container_width=True)
    else:
        with st.expander("📕 PDF Not Available"):
            st.code("pip install reportlab")
            st.info("Install ReportLab to enable PDF export")

# ======================================
# TRANSLATION HELPER (for UI switching)
# ======================================
def translate_text(client, text: str, target_lang: str) -> str:
    """
    Translate given text to target_lang ('ar' or 'en').
    Uses caching in session_state to avoid repeated API calls.
    """
    if not text:
        return ""

    cache_key = f"translation_{target_lang}_" + str(hash(text))
    if cache_key in st.session_state:
        return st.session_state[cache_key]

    try:
        if target_lang == "ar":
            prompt = f"Translate the following text to clear Arabic suitable for students:\n\n{text}"
        else:
            prompt = f"Translate the following text to clear English suitable for students:\n\n{text}"

        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
            max_tokens=2000
        )
        translated = response.choices[0].message.content
        st.session_state[cache_key] = translated
        return translated
    except Exception as e:
        return f"❌ Translation error: {str(e)}"


# ======================================
# AI FEATURES (أحمد) - AR + EN
# ======================================

def generate_summary(client, text: str, lang: str) -> str:
    """Generate summary in given lang ('en' or 'ar')."""
    try:
        if lang == 'ar':
            instructions = """
اكتب ملخصًا منظمًا باللغة العربية:
- قسم الملخص إلى عناوين رئيسية فرعية.
- تحت كل عنوان، اكتب نقاط مختصرة وواضحة.
- ركّز على التعريفات، القوانين، والأفكار الأساسية.
- اجعل الملخص مناسبًا للمراجعة السريعة قبل الامتحان.
"""
        else:
            instructions = """
Write a well-structured summary in English:
- Divide the summary into clear sections with headings.
- Under each heading, use concise bullet points.
- Focus on definitions, key concepts, and main ideas.
- Make it suitable for quick exam revision.
"""

        prompt = f"""
You are an expert in summarizing academic content.
{instructions}

Text:
{text[:15000]}

Structured Summary:
"""
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
            max_tokens=2000
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"❌ Error generating summary ({lang}): {str(e)}"


def generate_quiz(client, text: str, num_questions: int, lang: str) -> dict:
    """Generate quiz in given lang ('en' or 'ar')."""
    try:
        if lang == 'ar':
            lang_rules = """
- اكتب كل الأسئلة والاختيارات والتوضيحات باللغة العربية.
- اجعل صياغة الأسئلة بسيطة وواضحة للطلاب.
"""
        else:
            lang_rules = """
- Write all questions, options, and explanations in clear English.
"""

        prompt = f"""
Create a quiz in PURE JSON format. No markdown, no explanation, ONLY JSON.
Return exactly this structure:
{{
  "questions": [
    {{
      "question": "What is the capital of France?",
      "options": ["London", "Paris", "Berlin", "Madrid"],
      "correct_answer": 1,
      "explanation": "Paris is the capital and largest city of France."
    }}
  ]
}}

Rules:
- {num_questions} questions total
- Each question has exactly 4 options
- correct_answer is the index (0, 1, 2, or 3)
- Cover different key topics from the text
{lang_rules}
- Return ONLY valid JSON (no extra text)

Text to create quiz from:
{text[:15000]}

JSON Quiz:
"""
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=3000
        )

        content = response.choices[0].message.content.strip()
        content = content.replace('```json', '').replace('```', '').strip()

        json_match = re.search(r'\{.*\}', content, re.DOTALL)
        if json_match:
            content = json_match.group(0)

        quiz_data = json.loads(content)

        if 'questions' not in quiz_data or not isinstance(quiz_data['questions'], list):
            return {"questions": []}

        valid_questions = []
        for q in quiz_data['questions']:
            if (
                isinstance(q, dict)
                and 'question' in q
                and 'options' in q
                and 'correct_answer' in q
                and isinstance(q['options'], list)
                and len(q['options']) == 4
                and isinstance(q['correct_answer'], int)
                and 0 <= q['correct_answer'] < 4
            ):
                valid_questions.append(q)

        return {"questions": valid_questions}
    except Exception as e:
        return {"questions": [], "error": f"❌ Quiz error ({lang}): {str(e)}"}


def generate_mindmap(client, text: str, lang: str) -> str:
    """Generate mind map in given lang ('en' or 'ar')."""
    try:
        if lang == 'ar':
            instructions = """
أنت خبير في إنشاء خرائط ذهنية تعليمية.
أنشئ خريطة ذهنية منظمة من النص التالي.

المتطلبات:
- اجعل الفكرة الرئيسية واضحة في البداية.
- أنشئ من 4 إلى 6 فروع رئيسية فقط.
- كل فرع رئيسي يحتوي على 2 إلى 5 نقاط فرعية.
- استخدم الإيموجي بشكل بسيط لجعل الشكل جذابًا.
- استخدم تنسيقًا هرميًا واضحًا (عناوين رئيسية ثم نقاط فرعية).
- اكتب كل شيء باللغة العربية.
"""
        else:
            instructions = """
You are an expert in creating stunning educational mind maps.
Create a beautifully organized mind map from the following text.

Requirements:
- Start with a clear, concise MAIN TOPIC.
- Create 4–6 MAIN BRANCHES only.
- Each main branch should have 2–5 sub-points.
- Use emojis sparingly for visual appeal.
- Use clear hierarchical formatting (main topic → branches → sub-points).
- Write everything in English.
"""

        prompt = f"""
{instructions}

Text:
{text[:15000]}

Mind Map:
"""
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            max_tokens=2500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"❌ Error generating mind map ({lang}): {str(e)}"


def generate_question_bank(client, text: str, num_questions: int, lang: str) -> str:
    """Generate question bank in given lang ('en' or 'ar')."""
    try:
        if lang == 'ar':
            instructions = f"""
أنت خبير في إعداد بنوك الأسئلة للامتحانات.

أنشئ بنك أسئلة احترافي من النص التالي.

المتطلبات:
- إجمالي عدد الأسئلة: {num_questions}
- التوزيع:
  - صح/خطأ (~30%)
  - اختيار من متعدد (~40%)
  - أسئلة مقالية/قصيرة (~30%)
- ضع عنوانًا واضحًا لكل قسم:
  - أولًا: أسئلة صح وخطأ
  - ثانيًا: أسئلة اختيار من متعدد
  - ثالثًا: أسئلة مقالية قصيرة
- اجعل المسافات بين الأسئلة واضحة.
- إجابات أسئلة الصح والخطأ بين قوسين مثل: (صح) أو (خطأ).
- اجعل اللغة عربية واضحة وبسيطة.
"""
        else:
            instructions = f"""
You are an expert exam creator.

Create a professional, well-formatted question bank from the text below.

Requirements:
- Total questions: {num_questions}
- Distribution:
  - True/False (~30%)
  - Multiple Choice (~40%)
  - Short Answer (~30%)
- Add clear section headings:
  - Section 1: True/False Questions
  - Section 2: Multiple Choice Questions
  - Section 3: Short Answer Questions
- Use clear spacing between questions.
- True/False answers MUST be in parentheses like (True) or (False).
- Provide brief explanations or notes after the answer when possible.
- Use clear, exam-style English.
"""

        prompt = f"""
{instructions}

Text:
{text[:15000]}

Question Bank:
"""
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=4500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"❌ Error generating question bank ({lang}): {str(e)}"


def generate_all_features(client, text: str, filename: str):
    """
    Generate all 4 features automatically after file upload
    in BOTH English and Arabic, once.
    """
    # Use progress indicators
    progress_bar = st.progress(0)
    status_text = st.empty()

    features_en = {}
    features_ar = {}

    try:
        # ---------- ENGLISH ----------
        status_text.text("🤖 Generating English Summary...")
        progress_bar.progress(5)
        features_en['summary'] = generate_summary(client, text, 'en')

        status_text.text("🤖 Generating English Quiz...")
        progress_bar.progress(20)
        features_en['quiz'] = generate_quiz(client, text, num_questions=10, lang='en')

        status_text.text("🤖 Generating English Mind Map...")
        progress_bar.progress(35)
        features_en['mindmap'] = generate_mindmap(client, text, 'en')

        status_text.text("🤖 Generating English Question Bank...")
        progress_bar.progress(50)
        features_en['questionbank'] = generate_question_bank(client, text, num_questions=20, lang='en')

        # ---------- ARABIC ----------
        status_text.text("🤖 جاري إنشاء الملخص العربي...")
        progress_bar.progress(60)
        features_ar['summary'] = generate_summary(client, text, 'ar')

        status_text.text("🤖 جاري إنشاء الكويز العربي...")
        progress_bar.progress(75)
        features_ar['quiz'] = generate_quiz(client, text, num_questions=10, lang='ar')

        status_text.text("🤖 جاري إنشاء الخريطة الذهنية بالعربي...")
        progress_bar.progress(85)
        features_ar['mindmap'] = generate_mindmap(client, text, 'ar')

        status_text.text("🤖 جاري إنشاء بنك الأسئلة بالعربي...")
        progress_bar.progress(95)
        features_ar['questionbank'] = generate_question_bank(client, text, num_questions=20, lang='ar')

        progress_bar.progress(100)
        status_text.success("✅ All Arabic & English features generated!")
        time.sleep(1)
        status_text.empty()
        progress_bar.empty()

        # Store in session state
        st.session_state['features_en'] = features_en
        st.session_state['features_ar'] = features_ar
        st.session_state['last_filename'] = filename
        st.session_state['text_hash'] = get_text_hash(text)
        st.session_state['features_generated'] = True

        return {"en": features_en, "ar": features_ar}
    except Exception as e:
        status_text.error(f"❌ Error during generation: {str(e)}")
        return {"en": features_en, "ar": features_ar}

# ======================================
# RENDER HELPER FOR ARABIC HTML - بدون خلفية بلون أبيض
# ======================================
def render_content(text: str):
    """Render AI text nicely in Arabic (RTL) or English (Markdown)."""
    lang = st.session_state.get('ui_lang', 'ar')
    
    if lang == 'ar':
        # تنظيف النص من أي تنسيقات HTML أو Markdown
        import re
        
        # إزالة أي وسوم HTML موجودة
        text = re.sub(r'<[^>]+>', '', text)
        
        # إزالة تنسيقات Markdown
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  # **نص**
        text = re.sub(r'\*(.*?)\*', r'\1', text)      # *نص*
        text = re.sub(r'__(.*?)__', r'\1', text)      # __نص__
        text = re.sub(r'_(.*?)_', r'\1', text)        # _نص_
        text = re.sub(r'`(.*?)`', r'\1', text)        # `نص`
        
        # عرض النص بلون أبيض (بدون خلفية)
        st.markdown(f"""
        <div style="
            direction: rtl;
            text-align: right;
            padding: 10px;
            font-family: 'Segoe UI', 'Arial', sans-serif;
            font-size: 16px;
            line-height: 1.8;
            white-space: pre-wrap;
            word-wrap: break-word;
            color: #FFFFFF;
        ">
        {text}
        </div>
        """, unsafe_allow_html=True)
    else:
        # الإنجليزية - عرض عادي
        st.markdown(text)


# ======================================
# INTERACTIVE QUIZ DISPLAY (RTL READY)
# ======================================
def display_interactive_quiz(quiz_data: dict, filename: str):
    """Display interactive quiz with immediate feedback (RTL friendly)."""
    lang = st.session_state.get('ui_lang', 'ar')
    is_ar = (lang == 'ar')

    if not quiz_data or 'questions' not in quiz_data or len(quiz_data['questions']) == 0:
        st.error(T("❌ لا توجد أسئلة كويز متاحة", "❌ No quiz questions available"))
        st.info(T(
            "💡 ممكن الـ AI واجه مشكلة في إنشاء الكويز. جرّب ترفع ملف تاني أو تعيد الرفع.",
            "💡 The AI might have had trouble generating the quiz. Try uploading the file again or use a different document."
        ))
        return

    if is_ar:
        st.markdown("""
            <style>
            div.stRadio > div {
                direction: rtl;
                text-align: right;
            }
            div.stRadio label {
                direction: rtl;
                text-align: right;
            }
            </style>
        """, unsafe_allow_html=True)

    st.subheader(T("❓ كويز تفاعلي", "❓ Interactive Quiz"))
    st.markdown(T(
        f"**عدد الأسئلة: {len(quiz_data['questions'])} - جاوب وخد فيدباك فورًا!**",
        f"**{len(quiz_data['questions'])} Questions - Answer each question and get instant feedback!**"
    ))
    st.markdown("---")

    if 'quiz_answers' not in st.session_state:
        st.session_state['quiz_answers'] = {}
    if 'quiz_submitted' not in st.session_state:
        st.session_state['quiz_submitted'] = {}

    total_questions = len(quiz_data['questions'])
    correct_count = 0

    for idx, q in enumerate(quiz_data['questions']):
        question_key = f"q_{idx}"

        st.markdown(T(
            f"### السؤال {idx + 1} من {total_questions}",
            f"### Question {idx + 1} of {total_questions}"
        ))
        st.markdown(f"**{q['question']}**")

        user_answer = st.radio(
            T("اختر الإجابة:", "Select your answer:"),
            options=range(len(q['options'])),
            format_func=lambda x: f"{chr(65+x)}) {q['options'][x]}",
            key=f"radio_{filename}_{question_key}",
            index=None,
            horizontal=False
        )

        col1, col2 = st.columns([1, 4])
        with col1:
            submit_btn = st.button(
                T("تأكيد", "Submit"),
                key=f"submit_{filename}_{question_key}"
            )

        if submit_btn and user_answer is not None:
            st.session_state['quiz_answers'][question_key] = user_answer
            st.session_state['quiz_submitted'][question_key] = True

        if st.session_state['quiz_submitted'].get(question_key, False):
            user_ans = st.session_state['quiz_answers'][question_key]
            correct_ans = q['correct_answer']

            if user_ans == correct_ans:
                st.success(T("✅ إجابة صحيحة", "✅ Correct!"))
                correct_count += 1
            else:
                st.error(T(
                    f"❌ إجابة خاطئة. الإجابة الصحيحة: {chr(65+correct_ans)}) {q['options'][correct_ans]}",
                    f"❌ Wrong! The correct answer is: {chr(65+correct_ans)}) {q['options'][correct_ans]}"
                ))

            if 'explanation' in q and q['explanation']:
                st.info(T(
                    f"💡 التوضيح: {q['explanation']}",
                    f"💡 Explanation: {q['explanation']}"
                ))

        st.markdown("---")

    submitted_count = sum(1 for k in st.session_state['quiz_submitted'].values() if k)
    if submitted_count == total_questions:
        score_percentage = (correct_count / total_questions) * 100
        st.markdown(T("### 📊 نتيجتك", "### 📊 Your Score"))
        st.metric(
            T("النتيجة", "Score"),
            f"{correct_count}/{total_questions} ({score_percentage:.1f}%)"
        )

        if score_percentage >= 80:
            st.balloons()
            st.success(T("🎉 ممتاز! شغل عالي!", "🎉 Excellent work!"))
        elif score_percentage >= 60:
            st.info(T("👍 جيد! استمر في المذاكرة.", "👍 Good job! Keep practicing!"))
        else:
            st.warning(T("📚 محتاج تراجع المحتوى وتعيد المحاولة.", "📚 Review the material and try again!"))

        if st.button(T("🔄 إعادة الكويز", "🔄 Reset Quiz")):
            st.session_state['quiz_answers'] = {}
            st.session_state['quiz_submitted'] = {}
            st.rerun()


# ======================================
# STREAMLIT UI - MAIN
# ======================================
def main():
    if 'ui_lang' not in st.session_state:
        st.session_state['ui_lang'] = 'ar'
    
    # Initialize session state variables if they don't exist
    if 'features_en' not in st.session_state:
        st.session_state['features_en'] = {}
    if 'features_ar' not in st.session_state:
        st.session_state['features_ar'] = {}
    if 'last_filename' not in st.session_state:
        st.session_state['last_filename'] = None
    if 'last_file_hash' not in st.session_state:
        st.session_state['last_file_hash'] = None
    if 'text_hash' not in st.session_state:
        st.session_state['text_hash'] = None
    if 'selected_feature' not in st.session_state:
        st.session_state['selected_feature'] = 'summarize'
    if 'features_generated' not in st.session_state:
        st.session_state['features_generated'] = False
    if 'processed_files' not in st.session_state:
        st.session_state['processed_files'] = {}  # تخزين الملفات التي تمت معالجتها

    with st.sidebar:
        st.header("⚙️ Settings")
        lang = st.radio(
            "Language / اللغة",
            options=['ar', 'en'],
            format_func=lambda x: "🇪🇬 العربية" if x == 'ar' else "🇬🇧 English",
            key="ui_lang_radio"
        )
        st.session_state['ui_lang'] = lang

        api_key = st.text_input(
            T("🔑 مفتاح Mistral API", "🔑 Mistral API Key"),
            type="password",
            help=T("اكتب مفتاح الـ API الخاص بـ Mistral", "Enter your Mistral API key")
        )

        if not api_key:
            st.warning(T(
                "⚠️ من فضلك أدخل مفتاح الـ API للبدء",
                "⚠️ Please enter your API Key to start"
            ))
            st.info("Get your free API key from: https://console.mistral.ai/")
            st.stop()

        client = cache.get_mistral_client(api_key)
        if not client:
            st.error(T(
                "❌ فشل الاتصال بـ Mistral API",
                "❌ Failed to connect to Mistral API"
            ))
            st.stop()
        st.success(T("✅ تم الاتصال بنجاح!", "✅ Connected successfully!"))
        st.markdown("---")

        st.subheader(T("📝 إعدادات المعالجة", "📝 Processing Settings"))
        chunk_size = st.slider(
            T("حجم الجزء (Chunk Size)", "Chunk Size"),
            min_value=300,
            max_value=1000,
            value=500,
            step=50,
            help=T(
                "حجم كل جزء من النص عند التقسيم",
                "Size of each text chunk"
            )
        )
        chunk_overlap = st.slider(
            T("تداخل الأجزاء (Overlap)", "Chunk Overlap"),
            min_value=50,
            max_value=300,
            value=100,
            step=25,
            help=T(
                "مقدار التداخل بين الأجزاء",
                "Overlap between chunks"
            )
        )

        st.markdown("---")
        st.info(T(
            "💡 **كل المميزات هتتولد تلقائي بعد رفع الملف!**",
            "💡 **All features are generated automatically after file upload!**"
        ))

        if not PDF_AVAILABLE:
            st.warning("⚠️ PDF export disabled. Install: `pip install reportlab`")
        else:
            st.success("✅ PDF export enabled")

    lang = st.session_state.get('ui_lang', 'ar')
    direction = "rtl" if lang == 'ar' else "ltr"
    align = "right" if lang == 'ar' else "left"

    st.markdown(f"""
        <style>
        .main, .block-container {{
            direction: {direction};
            text-align: {align};
        }}
        div[role="textbox"], textarea, input, label, p, span, h1, h2, h3, h4, h5 {{
            direction: {direction};
            text-align: {align};
        }}
        /* تنسيق الأزرار لتكون بنفس الحجم */
        div.stButton > button {{
            width: 100%;
            height: 50px;
            font-size: 16px;
            font-weight: bold;
        }}
        </style>
    """, unsafe_allow_html=True)

    st.title(T("🎓 مساعد الدراسة بالذكاء الاصطناعي", "🎓 AI Study Assistant"))
    st.markdown(T(
        "**ارفع ملفك وخد ملخص، كويز، خريطة ذهنية، وبنك أسئلة تلقائيًا**",
        "**Upload your document and get instant AI-powered study materials**"
    ))
    st.markdown("---")

    st.header(T("📤 رفع المستند", "📤 Upload Document"))
    uploaded_file = st.file_uploader(
        T("اختر ملف (PDF, Word, PowerPoint, Image, TXT)", "Choose a file (PDF, Word, PowerPoint, Image, TXT)"),
        type=["pdf", "docx", "doc", "pptx", "ppt", "txt", "png", "jpg", "jpeg"],
        help=T("الحد الأقصى للحجم: 50 ميجابايت", "Maximum file size: 50 MB")
    )

    if uploaded_file:
        is_valid, validation_msg = validate_file(uploaded_file)
        if not is_valid:
            st.error(validation_msg)
            st.stop()
        st.success(validation_msg)

        # Get file content and hash
        file_content = uploaded_file.read()
        uploaded_file.seek(0)
        current_hash = cache.get_file_hash(file_content)
        filename = uploaded_file.name
        
        # التحقق من أن الملف لم تتم معالجته من قبل
        is_new_file = current_hash not in st.session_state['processed_files']

        if is_new_file:
            # Reset features_generated flag for new file
            st.session_state['features_generated'] = False
            
            with st.spinner(T("🔄 جاري معالجة المستند...", "🔄 Processing document...")):
                ext = get_file_extension(filename)

                if ext == ".docx":
                    extracted_text, metadata = extract_from_word(uploaded_file)
                elif ext == ".txt":
                    extracted_text, metadata = extract_from_txt(uploaded_file)
                elif ext == ".pptx":
                    extracted_text, metadata = pptx_extract_text(uploaded_file, client)
                elif ext in [".pdf", ".png", ".jpg", ".jpeg"]:
                    extracted_text, metadata = ocr_mistral(uploaded_file, filename, client)
                else:
                    st.error(T("❌ نوع الملف غير مدعوم", "❌ Unsupported file type"))
                    st.stop()

                if extracted_text.startswith("Error") or extracted_text.startswith("❌"):
                    st.error(extracted_text)
                    st.stop()

                # تخزين النص المستخرج والميتاداتا
                st.session_state['extracted_text'] = extracted_text
                st.session_state['metadata'] = metadata
                st.session_state['filename'] = filename
                st.session_state['text_hash'] = get_text_hash(extracted_text)
                
                # تسجيل الملف كـ "تمت معالجته"
                st.session_state['processed_files'][current_hash] = {
                    'filename': filename,
                    'extracted_text': extracted_text,
                    'metadata': metadata,
                    'text_hash': st.session_state['text_hash']
                }

                st.success(T("✅ تم استخراج النص بنجاح!", "✅ Text extracted successfully!"))

                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric(T("📝 عدد الكلمات", "📝 Words"), f"{len(extracted_text.split()):,}")
                with col2:
                    st.metric(T("🔤 عدد الحروف", "🔤 Characters"), f"{len(extracted_text):,}")
                with col3:
                    chunks = chunk_text(extracted_text, chunk_size, chunk_overlap)
                    st.metric(T("📦 عدد الأجزاء", "📦 Chunks"), len(chunks))

                with st.expander(T("👀 عرض النص المستخرج", "👀 View Extracted Text")):
                    st.text_area(
                        T("النص الكامل", "Full Text"),
                        extracted_text,
                        height=300,
                        key="extracted_text_display"
                    )

                if PDF_AVAILABLE:
                    try:
                        pdf_buffer = create_pdf_from_text(
                            extracted_text,
                            T(f"النص المستخرج - {filename}", f"Extracted Text - {filename}")
                        )
                        st.download_button(
                            label=T("📥 تحميل النص المستخرج (PDF)", "📥 Download Extracted Text (PDF)"),
                            data=pdf_buffer,
                            file_name=f"{filename}_extracted.pdf",
                            mime="application/pdf",
                            key="download_extracted_pdf"
                        )
                    except Exception as e:
                        st.warning(f"⚠️ PDF Error: {str(e)[:50]}")
                        st.download_button(
                            label=T("📥 تحميل النص المستخرج (TXT)", "📥 Download Extracted Text (TXT)"),
                            data=extracted_text,
                            file_name=f"{filename}_extracted.txt",
                            mime="text/plain",
                            key="download_extracted_txt"
                        )
                else:
                    st.download_button(
                        label=T("📥 تحميل النص المستخرج (TXT)", "📥 Download Extracted Text (TXT)"),
                        data=extracted_text,
                        file_name=f"{filename}_extracted.txt",
                        mime="text/plain",
                        key="download_extracted_txt"
                    )

                st.markdown("---")

                # Generate features only for new files
                st.header(T("🤖 جاري إنشاء مواد المذاكرة...", "🤖 Generating AI Features..."))
                st.info(T(
                    "من فضلك انتظر لحظات حتى يتم إنشاء المحتوى بالعربي والإنجليزي.",
                    "Please wait while we generate Arabic and English study materials..."
                ))
                _ = generate_all_features(client, extracted_text, filename)
                st.success(T("✅ تم إنشاء النسخة العربية والإنجليزية!", "✅ Arabic & English versions generated!"))
                
                # Set default selected feature
                st.session_state['selected_feature'] = 'summarize'
                
        else:
            # نفس الملف - استرجاع البيانات المخزنة
            st.success(T("✅ تم استخراج النص بنجاح!", "✅ Text extracted successfully!"))
            
            # استرجاع البيانات من الذاكرة
            file_data = st.session_state['processed_files'][current_hash]
            st.session_state['extracted_text'] = file_data['extracted_text']
            st.session_state['metadata'] = file_data['metadata']
            st.session_state['filename'] = file_data['filename']
            st.session_state['text_hash'] = file_data['text_hash']
            
            # التأكد أن features_generated = True
            st.session_state['features_generated'] = True
            
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric(T("📝 عدد الكلمات", "📝 Words"), f"{len(st.session_state['extracted_text'].split()):,}")
            with col2:
                st.metric(T("🔤 عدد الحروف", "🔤 Characters"), f"{len(st.session_state['extracted_text']):,}")
            with col3:
                chunks = chunk_text(st.session_state['extracted_text'], chunk_size, chunk_overlap)
                st.metric(T("📦 عدد الأجزاء", "📦 Chunks"), len(chunks))

            with st.expander(T("👀 عرض النص المستخرج", "👀 View Extracted Text")):
                st.text_area(
                    T("النص الكامل", "Full Text"),
                    st.session_state['extracted_text'],
                    height=300,
                    key="extracted_text_display"
                )

    if 'extracted_text' in st.session_state and st.session_state['features_generated']:
        filename = st.session_state['filename']
        ui_lang = st.session_state.get('ui_lang', 'ar')

        st.header(T("🎯 مميزات المذاكرة", "🎯 Study Features"))
        st.markdown(T(
            "... **اضغط على أي زر لعرض المحتوى الناتج:**",
            "... **Click any button to view the generated content:**"
        ))

        # أزرار متساوية الحجم (5 أزرار في صف واحد)
        col1, col2, col3, col4, col5 = st.columns(5)
        
        with col1:
            btn_style = "primary" if st.session_state.get('selected_feature') == 'summarize' else "secondary"
            if st.button(T("📄 الملخص", "📄 Summary"), key="btn_summary", use_container_width=True, type=btn_style):
                st.session_state['selected_feature'] = 'summarize'
                st.rerun()
        
        with col2:
            btn_style = "primary" if st.session_state.get('selected_feature') == 'quiz' else "secondary"
            if st.button(T("❓ الكويز", "❓ Quiz"), key="btn_quiz", use_container_width=True, type=btn_style):
                st.session_state['selected_feature'] = 'quiz'
                st.rerun()
        
        with col3:
            btn_style = "primary" if st.session_state.get('selected_feature') == 'mindmap' else "secondary"
            if st.button(T("🧠 الخريطة الذهنية", "🧠 Mind Map"), key="btn_mindmap", use_container_width=True, type=btn_style):
                st.session_state['selected_feature'] = 'mindmap'
                st.rerun()
        
        with col4:
            btn_style = "primary" if st.session_state.get('selected_feature') == 'questionbank' else "secondary"
            if st.button(T("📚 بنك الأسئلة", "📚 Question Bank"), key="btn_qbank", use_container_width=True, type=btn_style):
                st.session_state['selected_feature'] = 'questionbank'
                st.rerun()
        
        with col5:
            btn_style = "primary" if st.session_state.get('selected_feature') == 'chatbot' else "secondary"
            if st.button(T("💬 المحادثة", "💬 Chat"), key="btn_chat", use_container_width=True, type=btn_style):
                st.session_state['selected_feature'] = 'chatbot'
                st.rerun()

        st.markdown("---")

        features_en = st.session_state.get('features_en', {})
        features_ar = st.session_state.get('features_ar', {})
        
        # اختيار النسخة حسب لغة الواجهة
        features = features_ar if ui_lang == 'ar' else features_en

        # عرض المحتوى حسب الزر المختار
        if st.session_state['selected_feature'] == 'summarize':
            st.subheader(T("📄 الملخص الذكي", "📄 Smart Summary"))
            summary_text = features.get('summary', '')
            if not summary_text:
                st.error(T("❌ لا يوجد ملخص متاح", "❌ Summary not available"))
            else:
                render_content(summary_text)
                st.markdown("---")
                st.markdown(T("### 📥 تحميل الملخص", "### 📥 Download Summary"))
                create_download_button_pdf(
                    summary_text,
                    f"{filename}_summary_{ui_lang}",
                    T("الملخص", "Summary"),
                    "download_summary"
                )

        elif st.session_state['selected_feature'] == 'quiz':
            quiz_data = features.get('quiz', {})
            if not quiz_data or not quiz_data.get('questions'):
                st.error(T("❌ لا يوجد كويز متاح", "❌ Quiz not available"))
            else:
                display_interactive_quiz(quiz_data, filename + "_" + ui_lang)

        elif st.session_state['selected_feature'] == 'mindmap':
            st.subheader(T("🧠 الخريطة الذهنية", "🧠 Mind Map"))
            mindmap_text = features.get('mindmap', '')
            if not mindmap_text:
                st.error(T("❌ لا توجد خريطة ذهنية متاحة", "❌ Mind Map not available"))
            else:
                render_content(mindmap_text)
                st.markdown("---")
                st.markdown(T("### 📥 تحميل الخريطة الذهنية", "### 📥 Download Mind Map"))
                create_download_button_pdf(
                    mindmap_text,
                    f"{filename}_mindmap_{ui_lang}",
                    T("الخريطة الذهنية", "Mind Map"),
                    "download_mindmap"
                )

        elif st.session_state['selected_feature'] == 'questionbank':
            st.subheader(T("📚 بنك الأسئلة", "📚 Question Bank"))
            qb_text = features.get('questionbank', '')
            if not qb_text:
                st.error(T("❌ لا يوجد بنك أسئلة متاح", "❌ Question Bank not available"))
            else:
                render_content(qb_text)
                st.markdown("---")
                st.markdown(T("### 📥 تحميل بنك الأسئلة", "### 📥 Download Question Bank"))
                create_download_button_pdf(
                    qb_text,
                    f"{filename}_questionbank_{ui_lang}",
                    T("بنك الأسئلة", "Question Bank"),
                    "download_questionbank"
                )

        elif st.session_state['selected_feature'] == 'chatbot':
            display_chat_interface(client, filename)

if __name__ == "__main__":
    main()