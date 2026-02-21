"""
Utility functions for AI Study Assistant API
Extracted text processing, AI feature generation, and helper functions
"""

import os
import base64
import json
import re
import hashlib
from io import BytesIO
from datetime import datetime
from typing import Optional, List, Dict, Tuple

from docx import Document as DocxDocument
from pptx import Presentation

# LangChain replaced with built-in splitter for Python 3.14 compatibility
class RecursiveCharacterTextSplitter:
    def __init__(self, chunk_size=500, chunk_overlap=100, separators=None, length_function=len):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = separators or ["\n\n", "\n", ". ", " ", ""]
        self.length_function = length_function

    def split_text(self, text: str) -> List[str]:
        chunks = []
        start = 0
        while start < len(text):
            end = start + self.chunk_size
            chunks.append(text[start:end])
            start += self.chunk_size - self.chunk_overlap
        return [c for c in chunks if c.strip()]

# ======================================
# FILE VALIDATION & UTILITIES
# ======================================
def validate_file(file_content: bytes, filename: str, max_size_mb: float = 50) -> Tuple[bool, str]:
    """Validate file before processing."""
    try:
        if not file_content or len(file_content) == 0:
            return False, "❌ Empty file detected"
        size_mb = len(file_content) / (1024 * 1024)
        if size_mb > max_size_mb:
            return False, f"❌ File too large: {size_mb:.1f}MB (Max: {max_size_mb}MB)"
        return True, f"✅ Valid file ({size_mb:.2f}MB)"
    except Exception as e:
        return False, f"❌ Validation error: {str(e)}"

def get_file_extension(filename: str) -> str:
    """Get file extension safely."""
    return os.path.splitext(filename)[1].lower()

def get_text_hash(text: str) -> str:
    """Generate a hash for the text."""
    return hashlib.md5(text.encode('utf-8')).hexdigest()

# ======================================
# TEXT CLEANING & PREPROCESSING
# ======================================
def clean_text(text: str) -> str:
    """Clean the extracted text - Basic cleaning."""
    if not text:
        return ""
    text = text.replace("\t", " ")
    text = text.replace("\xa0", " ")
    text = " ".join(text.split())
    return text.strip()

def chunk_text(text: str, chunk_size: int = 500, chunk_overlap: int = 100) -> List[str]:
    """Split text into overlapping chunks."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
        length_function=len,
    )
    chunks = splitter.split_text(text)
    return chunks

# ======================================
# FILE EXTRACTION FUNCTIONS
# ======================================
def extract_from_word(file_obj) -> Tuple[str, dict]:
    """Extract text from Word document."""
    try:
        doc = DocxDocument(file_obj)
        all_text = []
        metadata = {
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
            "sections": len(doc.sections),
            "errors": []
        }

        for para in doc.paragraphs:
            if para.text.strip():
                all_text.append(para.text)

        for table in doc.tables:
            all_text.append("\n[Table Start]")
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    all_text.append(" | ".join(row_text))
            all_text.append("[Table End]\n")

        final_text = "\n\n".join(all_text)
        return clean_text(final_text), metadata

    except Exception as e:
        return f"Error extracting from Word: {str(e)}", {"error": str(e)}

def extract_from_txt(file_obj) -> Tuple[str, dict]:
    """Extract text from TXT file with multiple encoding support."""
    encodings = ['utf-8', 'utf-8-sig', 'cp1256', 'latin-1', 'iso-8859-1']
    for encoding in encodings:
        try:
            file_obj.seek(0)
            content = file_obj.read().decode(encoding)
            lines = content.split('\n')
            metadata = {
                "lines": len(lines),
                "characters": len(content),
                "words": len(content.split()),
                "encoding": encoding
            }
            return clean_text(content), metadata
        except (UnicodeDecodeError, AttributeError):
            continue

    return "Error: Could not decode text file with supported encodings", {
        "error": "Encoding not supported",
        "tried_encodings": encodings
    }

def pptx_extract_text(file_obj, client) -> Tuple[str, dict]:
    """Extract text + tables + OCR from images inside PPTX."""
    try:
        prs = Presentation(file_obj)
        all_slides_output = []
        metadata = {
            "total_slides": len(prs.slides),
            "slides_with_images": 0,
            "slides_with_tables": 0,
            "total_images_ocr": 0,
            "successful_ocr": 0,
            "failed_ocr": 0,
            "total_tables": 0,
            "errors": []
        }

        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_text_parts = [f"=== Slide {slide_number} ==="]
            has_images = False
            has_tables = False

            # Extract text
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            slide_text_parts.append(text)
                    elif hasattr(shape, "text") and isinstance(shape.text, str):
                        text = shape.text.strip()
                        if text:
                            slide_text_parts.append(text)
                except Exception as e:
                    metadata["errors"].append(f"Slide {slide_number}: Text extraction error - {str(e)}")

            # Extract tables
            for shape in slide.shapes:
                try:
                    if shape.has_table:
                        has_tables = True
                        metadata["total_tables"] += 1
                        table = shape.table
                        slide_text_parts.append("\n[Table Start]")
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_data:
                                slide_text_parts.append(" | ".join(row_data))
                        slide_text_parts.append("[Table End]\n")
                except Exception as e:
                    metadata["errors"].append(f"Slide {slide_number}: Table extraction error - {str(e)}")

            # Extract images with OCR
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "image"):
                    has_images = True
                    metadata["total_images_ocr"] += 1
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        mime = image.ext if image.ext in ["png", "jpg", "jpeg"] else "png"

                        image_size_mb = len(image_bytes) / (1024 * 1024)
                        if image_size_mb > 10:
                            metadata["errors"].append(f"Slide {slide_number}, Image {shape_idx}: Image too large")
                            metadata["failed_ocr"] += 1
                            continue

                        encoded = base64.b64encode(image_bytes).decode("utf-8")
                        document = {
                            "type": "image_url",
                            "image_url": f"data:image/{mime};base64,{encoded}"
                        }

                        response = client.ocr.process(
                            model="mistral-ocr-latest",
                            document=document,
                            include_image_base64=False
                        )

                        if not hasattr(response, 'pages') or not response.pages:
                            metadata["failed_ocr"] += 1
                            continue

                        ocr_text_found = False
                        for page in response.pages:
                            page_text = getattr(page, 'markdown', None) or getattr(page, 'text', None)
                            if page_text:
                                slide_text_parts.append(f"\n[Image OCR - Shape {shape_idx}]\n{page_text}")
                                ocr_text_found = True
                                metadata["successful_ocr"] += 1

                        if not ocr_text_found:
                            metadata["failed_ocr"] += 1
                    except Exception as e:
                        metadata["errors"].append(f"Slide {slide_number}, Image {shape_idx}: {str(e)}")
                        metadata["failed_ocr"] += 1

            if has_images:
                metadata["slides_with_images"] += 1
            if has_tables:
                metadata["slides_with_tables"] += 1

            slide_output = "\n".join(slide_text_parts)
            all_slides_output.append(slide_output)

        final_text = "\n\n--- Slide Break ---\n\n".join(all_slides_output)
        return final_text.strip(), metadata

    except Exception as e:
        return f"Error processing PPTX: {str(e)}", {"error": str(e)}

def ocr_mistral(file_obj, filename: str, client) -> Tuple[str, dict]:
    """OCR for images/PDF using Mistral."""
    metadata = {
        "filename": filename,
        "pages_processed": 0,
        "errors": []
    }
    try:
        file_bytes = file_obj.read()
        file_obj.seek(0)
        file_size_mb = len(file_bytes) / (1024 * 1024)
        metadata["file_size_mb"] = round(file_size_mb, 2)

        if file_size_mb > 50:
            error_msg = f"❌ File exceeds Mistral OCR limit of 50MB"
            metadata["errors"].append(error_msg)
            return error_msg, metadata

        if len(file_bytes) == 0:
            return "Error: Empty file", metadata

        encoded = base64.b64encode(file_bytes).decode("utf-8")
        ext = get_file_extension(filename)

        if ext == ".pdf":
            document = {
                "type": "document_url",
                "document_url": f"data:application/pdf;base64,{encoded}"
            }
        else:
            document = {
                "type": "image_url",
                "image_url": f"data:image/{ext.replace('.', '')};base64,{encoded}"
            }

        response = client.ocr.process(
            model="mistral-ocr-latest",
            document=document,
            include_image_base64=False
        )

        final_text = ""
        pages = getattr(response, "pages", [])
        metadata["pages_processed"] = len(pages)

        for page in pages:
            page_text = getattr(page, "markdown", None) or getattr(page, "text", "")
            final_text += page_text + "\n\n"

        return final_text.strip(), metadata

    except Exception as e:
        error_msg = f"Error processing file: {str(e)}"
        metadata["errors"].append(error_msg)
        return error_msg, metadata

# ======================================
# AI FEATURE GENERATION
# ======================================
def generate_summary(client, text: str, lang: str) -> str:
    """Generate summary in given lang ('en' or 'ar')."""
    try:
        if lang == 'ar':
            instructions = """
اكتب ملخصًا منظمًا باللغة العربية:
- قسم الملخص إلى عناوين رئيسية فرعية.
- تحت كل عنوان، اكتب نقاط مختصرة وواضحة.
- ركّز على التعريفات، القوانين، والأفكار الأساسية.
- اجعل الملخص مناسبًا للمراجعة السريعة قبل الامتحان.
"""
        else:
            instructions = """
Write a well-structured summary in English:
- Divide the summary into clear sections with headings.
- Under each heading, use concise bullet points.
- Focus on definitions, key concepts, and main ideas.
- Make it suitable for quick exam revision.
"""

        prompt = f"""
You are an expert in summarizing academic content.
{instructions}

Text:
{text[:15000]}

Structured Summary:
"""
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
            max_tokens=2000
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"❌ Error generating summary ({lang}): {str(e)}"

def generate_quiz(client, text: str, num_questions: int, lang: str) -> dict:
    """Generate quiz in given lang ('en' or 'ar')."""
    try:
        if lang == 'ar':
            lang_rules = """
- اكتب كل الأسئلة والاختيارات والتوضيحات باللغة العربية.
- اجعل صياغة الأسئلة بسيطة وواضحة للطلاب.
"""
        else:
            lang_rules = """
- Write all questions, options, and explanations in clear English.
"""

        prompt = f"""
Create a quiz in PURE JSON format. No markdown, no explanation, ONLY JSON.
Return exactly this structure:
{{
  "questions": [
    {{
      "question": "What is the capital of France?",
      "options": ["London", "Paris", "Berlin", "Madrid"],
      "correct_answer": 1,
      "explanation": "Paris is the capital and largest city of France."
    }}
  ]
}}

Rules:
- {num_questions} questions total
- Each question has exactly 4 options
- correct_answer is the index (0, 1, 2, or 3)
- Cover different key topics from the text
{lang_rules}
- Return ONLY valid JSON (no extra text)

Text to create quiz from:
{text[:15000]}

JSON Quiz:
"""
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=3000
        )

        content = response.choices[0].message.content.strip()
        content = content.replace('```json', '').replace('```', '').strip()

        json_match = re.search(r'\{.*\}', content, re.DOTALL)
        if json_match:
            content = json_match.group(0)

        quiz_data = json.loads(content)

        if 'questions' not in quiz_data or not isinstance(quiz_data['questions'], list):
            return {"questions": []}

        valid_questions = []
        for q in quiz_data['questions']:
            if (
                isinstance(q, dict)
                and 'question' in q
                and 'options' in q
                and 'correct_answer' in q
                and isinstance(q['options'], list)
                and len(q['options']) == 4
                and isinstance(q['correct_answer'], int)
                and 0 <= q['correct_answer'] < 4
            ):
                valid_questions.append(q)

        return {"questions": valid_questions}
    except Exception as e:
        return {"questions": [], "error": f"❌ Quiz error ({lang}): {str(e)}"}

def generate_mindmap(client, text: str, lang: str) -> str:
    """Generate mind map in given lang ('en' or 'ar')."""
    try:
        if lang == 'ar':
            instructions = """
أنت خبير في إنشاء خرائط ذهنية تعليمية.
أنشئ خريطة ذهنية منظمة من النص التالي.

المتطلبات:
- اجعل الفكرة الرئيسية واضحة في البداية.
- أنشئ من 4 إلى 6 فروع رئيسية فقط.
- كل فرع رئيسي يحتوي على 2 إلى 5 نقاط فرعية.
- استخدم الإيموجي بشكل بسيط لجعل الشكل جذابًا.
- استخدم تنسيقًا هرميًا واضحًا (عناوين رئيسية ثم نقاط فرعية).
- اكتب كل شيء باللغة العربية.
"""
        else:
            instructions = """
You are an expert in creating stunning educational mind maps.
Create a beautifully organized mind map from the following text.

Requirements:
- Start with a clear, concise MAIN TOPIC.
- Create 4-6 MAIN BRANCHES only.
- Each main branch should have 2-5 sub-points.
- Use emojis sparingly for visual appeal.
- Use clear hierarchical formatting (main topic → branches → sub-points).
- Write everything in English.
"""

        prompt = f"""
{instructions}

Text:
{text[:15000]}

Mind Map:
"""
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            max_tokens=2500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"❌ Error generating mind map ({lang}): {str(e)}"

def generate_question_bank(client, text: str, num_questions: int, lang: str) -> str:
    """Generate question bank in given lang ('en' or 'ar')."""
    try:
        if lang == 'ar':
            instructions = f"""
أنت خبير في إعداد بنوك الأسئلة للامتحانات.

أنشئ بنك أسئلة احترافي من النص التالي.

المتطلبات:
- إجمالي عدد الأسئلة: {num_questions}
- التوزيع:
  - صح/خطأ (~30%)
  - اختيار من متعدد (~40%)
  - أسئلة مقالية/قصيرة (~30%)
- ضع عنوانًا واضحًا لكل قسم.
- اجعل اللغة عربية واضحة وبسيطة.
"""
        else:
            instructions = f"""
You are an expert exam creator.

Create a professional, well-formatted question bank from the text below.

Requirements:
- Total questions: {num_questions}
- Distribution:
  - True/False (~30%)
  - Multiple Choice (~40%)
  - Short Answer (~30%)
- Add clear section headings.
- Use clear, exam-style English.
"""

        prompt = f"""
{instructions}

Text:
{text[:15000]}

Question Bank:
"""
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=4500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"❌ Error generating question bank ({lang}): {str(e)}"

def generate_chat_response(client, context: str, user_message: str, lang: str) -> str:
    """Generate AI response for chat based on document context."""
    try:
        if lang == 'ar':
            system_prompt = f"""أنت مساعد دراسة متخصص في شرح ومناقشة المستند المرفوع.

محتوى المستند:
{context[:8000]}

تعليمات مهمة:
1. يجب أن تجيب بناءً فقط على محتوى المستند أعلاه
2. إذا كان السؤال خارج نطاق المستند، قل بأدب أنك تستطيع فقط الإجابة بناءً على المستند المرفوع
3. كن مفيدًا وتعليميًا ودقيقًا
4. أجب باللغة العربية

سؤال المستخدم:
{user_message}

الإجابة:"""
        else:
            system_prompt = f"""You are an AI Study Assistant specialized in explaining and discussing the uploaded document.

DOCUMENT CONTEXT:
{context[:8000]}

IMPORTANT INSTRUCTIONS:
1. Answer based ONLY on the document content provided above
2. If the question is outside the document scope, politely say so
3. Be helpful, educational, and precise
4. Answer in English

USER'S QUESTION:
{user_message}

ANSWER:"""

        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message}
            ],
            temperature=0.7,
            max_tokens=1500
        )

        return response.choices[0].message.content

    except Exception as e:
        return f"❌ Error: {str(e)}"